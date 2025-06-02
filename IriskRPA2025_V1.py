
# step 2: generating the excel workbook and docx files

#pip install xlrd
#pip install openpyxl
#pip install pandas
#pip install docxtpl
# import openpyxl
# import xlrd


import openpyxl
import xlrd
import re
from openpyxl import Workbook
from openpyxl.styles import NamedStyle
from datetime import datetime

# Step 2.1: generating the excel workbook

# Define a custom date format style in openpyxl
date_style = NamedStyle(name="date_style", number_format="DD-MMM-YYYY")

# Open the .xls file using xlrd
xls_file_path = 'Export-Issue.xls'
workbook_xls = xlrd.open_workbook(xls_file_path)

# Create a new .xlsx workbook using openpyxl
workbook_xlsx = Workbook()
sheet_xlsx = workbook_xlsx.active

# Iterate through each sheet in the .xls file
for sheet_index in range(workbook_xls.nsheets):
    sheet_xls = workbook_xls.sheet_by_index(sheet_index)

    # Create a new sheet in the .xlsx workbook for each .xls sheet
    if sheet_index == 0:
        sheet_xlsx.title = sheet_xls.name
    else:
        sheet_xlsx = workbook_xlsx.create_sheet(title=sheet_xls.name)

        # Iterate through each row and column, copying the content to the .xlsx sheet
    for row in range(sheet_xls.nrows):
        for col in range(sheet_xls.ncols):
            cell_value = sheet_xls.cell_value(row, col)
            cell_type = sheet_xls.cell_type(row, col)

            # Check if the cell is a date
            if cell_type == xlrd.XL_CELL_DATE:
                date_value = xlrd.xldate.xldate_as_datetime(cell_value, workbook_xls.datemode)
                cell_xlsx = sheet_xlsx.cell(row=row + 1, column=col + 1, value=date_value)
                cell_xlsx.number_format = date_style.number_format
            else:
                sheet_xlsx.cell(row=row + 1, column=col + 1).value = cell_value

# Save the new .xlsx file
xlsx_file_path = 'Export-Issue.xlsx'
workbook_xlsx.save(xlsx_file_path)


from openpyxl import load_workbook

# Load the workbook and select the issue sheet
workbook = load_workbook('Export-Issue.xlsx')

# Delete File and Definition worksheets
worksheets_to_delete = ['Files', 'Definition']

# Delete the specified worksheets
for sheet_name in worksheets_to_delete:
    if sheet_name in workbook.sheetnames:  # Check if the sheet exists
        sheet = workbook[sheet_name]
        workbook.remove(sheet)

# Save the workbook after deleting the sheets
workbook.save('Export-Issue.xlsx')

ws = workbook.active


# Define the columns to keep (1-based index)
columns_to_keep = [7,8,16,24,27,30,33,39,40,42,45,62]

# Create a new workbook and sheet for the output
new_wb = Workbook()
new_ws = new_wb.active
new_ws.title = 'Selected Columns'

# Iterate through rows in the existing sheet and write the selected columns to the new sheet
for row in ws.iter_rows():
    new_row = [row[i-1].value for i in columns_to_keep]
    new_ws.append(new_row)

# Check if the cell is a date as dd-mmm-yyyy

# Define the desired date format style
date_style = NamedStyle(name="date_style", number_format="DD-MMM-YYYY")

# Function to check if a cell's value is a date
def is_date(value):
    if isinstance(value, datetime):
        return True
    return False

# Iterate over all cells in the sheet
for row in new_ws.iter_rows():
    for cell in row:
        if is_date(cell.value):
            # Check the current cell's number format
            if cell.number_format != "DD-MMM-YYYY":
                # Convert the cell's format to the desired format
                cell.number_format = date_style.number_format


# Function to check if any special mark and remove the special mark <> and replace this mark as ()
# Function to replace angle brackets <> with parentheses ()
def replace_special_marks(text):
    if isinstance(text, str):
        # Replace <...> with (...)
        return re.sub(r'<(.*?)>', r'(\1)', text)
    return text

# Apply the function to all cells in the new worksheet
for row in new_ws.iter_rows():
    for cell in row:
        cell.value = replace_special_marks(cell.value)


# Save the new workbook with the selected columns
new_wb.save('selected_columns.xlsx')

# Create two new sheets

accepted_ws = new_wb.create_sheet(title='Accepted')
open_ws = new_wb.create_sheet(title='Open')

# Assuming the first row contains headers, copy them to both new sheets
headers = [cell.value for cell in new_ws[1]]
accepted_ws.append(headers)
open_ws.append(headers)

# Iterate through the rows and distribute them based on the 'status' column
status_column = headers.index('Status') + 1  # Assuming the header row is 1-based

for row in new_ws.iter_rows(min_row=2, values_only=True):
    if row[status_column - 1] == 'Accepted':
        accepted_ws.append(row)
    elif row[status_column - 1] == 'Open':
        open_ws.append(row)
    elif row[status_column - 1] == 'Closed':
        open_ws.append(row)

# Iterate over all cells in the sheet
for row in accepted_ws.iter_rows():
    for cell in row:
        if is_date(cell.value):
            # Check the current cell's number format
            if cell.number_format != "DD-MMM-YYYY":
                # Convert the cell's format to the desired format
                cell.number_format = date_style.number_format

# Iterate over all cells in the sheet
for row in open_ws.iter_rows():
    for cell in row:
        if is_date(cell.value):
            # Check the current cell's number format
            if cell.number_format != "DD-MMM-YYYY":
                # Convert the cell's format to the desired format
                cell.number_format = date_style.number_format

new_wb.save('selected_columns.xlsx')


# Define the columns to in the accepted sheet to keep
new_ws = new_wb['Accepted']
columns_to_keep = ['Identifier','Issue Name','Current Issue Risk Rating','Review Date','Accepted date','Creation Date','Framework Process Indicator', 'Issue Description', 'Recommendation','Status Update']

# Find the indices of the columns to keep based on headers
header_row = [cell.value for cell in new_ws[1]]  # Assuming headers are in the first row
indices_to_keep = [header_row.index(col) + 1 for col in columns_to_keep]

# Create a list to store the rows with only the selected columns
filtered_rows = []

# Copy the header row for the selected columns
selected_headers = [header_row[i - 1] for i in indices_to_keep]
filtered_rows.append(selected_headers)

# Iterate through rows and copy the data for the selected columns
for row in new_ws.iter_rows(min_row=2, values_only=True):  # Start from row 2 to skip headers
    new_row = [row[i - 1] for i in indices_to_keep]
    filtered_rows.append(new_row)

# Clear the original worksheet content
new_ws.delete_rows(1, ws.max_row)

# Write the filtered data back to the original worksheet
for row in filtered_rows:
    new_ws.append(row)

# Iterate over all cells in the sheet
for row in accepted_ws.iter_rows():
    for cell in row:
        if is_date(cell.value):
            # Check the current cell's number format
            if cell.number_format != "DD-MMM-YYYY":
                # Convert the cell's format to the desired format
                cell.number_format = date_style.number_format

# Save the workbook with the modified worksheet
new_wb.save('selected_columns.xlsx')

# Part one Accepted risk finished


# Part two select High, Medium and low by Current Issue Risk Rating


# Determine Current Issue Risk Ranking

# Extract data from column 3 -Current Issue Risk Rating
from openpyxl import Workbook, load_workbook
import pandas as pd

new_wb = load_workbook('selected_columns.xlsx')
new_ws = new_wb['Open']

# Extract all data from the worksheet
data = []
for row in new_ws.iter_rows(values_only=True):
    data.append(row)

# Define the custom order for sorting
custom_order = {"High": 1, "Medium": 2, "Low": 3}

# Function to get the sorting key
def get_sort_key(value):
    return custom_order.get(value[2], 4)  # Default to 4 if the value is not found

# Sort data based on the third column (index 2)
sorted_data = sorted(data[1:], key=get_sort_key)  # Skip the header row
sorted_data.insert(0, data[0])  # Reinsert the header row at the top

# Clear the existing data
for row in new_ws.iter_rows():
    for cell in row:
        cell.value = None

# Write sorted data back to the worksheet
for i, row in enumerate(sorted_data, start=1):
    for j, value in enumerate(row, start=1):
        new_ws.cell(row=i, column=j, value=value)

# Iterate over all cells in the sheet
for row in open_ws.iter_rows():
    for cell in row:
        if is_date(cell.value):
            # Check the current cell's number format
            if cell.number_format != "DD-MMM-YYYY":
                # Convert the cell's format to the desired format
                cell.number_format = date_style.number_format

new_wb.save('selected_columns.xlsx')



# Create two new sheets -high findings or medium/low findings

High_ws = new_wb.create_sheet(title='High')
Medium_low_ws = new_wb.create_sheet(title='Medium_Low')

# Assuming the first row contains headers, copy them to both new sheets
headers = [cell.value for cell in new_ws[1]]
High_ws.append(headers)
Medium_low_ws.append(headers)

# Iterate through the rows and distribute them based on the 'Current Issue Risk Rating' column
status_column = headers.index('Current Issue Risk Rating') + 1  # Assuming the header row is 1-based
for row in new_ws.iter_rows(min_row=2, values_only=True):
    if row[status_column - 1] == 'High':
        High_ws.append(row)
    elif row[status_column - 1] in ['Medium', 'Low']:
        Medium_low_ws.append(row)

# Iterate over all cells in the sheet
for row in High_ws.iter_rows():
    for cell in row:
        if is_date(cell.value):
            # Check the current cell's number format
            if cell.number_format != "DD-MMM-YYYY":
                # Convert the cell's format to the desired format
                cell.number_format = date_style.number_format

# Iterate over all cells in the sheet
for row in Medium_low_ws.iter_rows():
    for cell in row:
        if is_date(cell.value):
            # Check the current cell's number format
            if cell.number_format != "DD-MMM-YYYY":
                # Convert the cell's format to the desired format
                cell.number_format = date_style.number_format

new_wb.save('selected_columns.xlsx')

# only keep the useful columns for high and medium worksheets

sheet1 = new_wb['High']
sheet2 = new_wb['Medium_Low']

# Define the columns to keep for each sheet

columns_to_keep_sheet1 = ['A', 'B', 'C', 'D','E', 'H','I', 'J', 'K', 'L']
columns_to_keep_sheet2 = ['A', 'B', 'C', 'D','E', 'H','I', 'J', 'K', 'L']

# Remove unwanted columns from the worksheet

# Convert column letters to indices
columns_to_keep_sheet1 = [openpyxl.utils.column_index_from_string(col) for col in columns_to_keep_sheet1]
columns_to_keep_sheet2 = [openpyxl.utils.column_index_from_string(col) for col in columns_to_keep_sheet2]

# Delete the unwanted columns from Sheet1
for col in range(sheet1.max_column, 0, -1):  # Loop backwards to avoid shifting issues
    if col not in columns_to_keep_sheet1:
        sheet1.delete_cols(col)

# Delete the unwanted columns from Sheet2
for col in range(sheet2.max_column, 0, -1):  # Loop backwards to avoid shifting issues
    if col not in columns_to_keep_sheet2:
        sheet2.delete_cols(col)


# arrange the column-creation date
from openpyxl import load_workbook
from openpyxl.utils import column_index_from_string

# Define the columns to move and the target positions
column_to_move_sheet1 = 'J'  # Column to move in Sheet1
column_to_move_sheet2 = 'J'  # Column to move in Sheet2
target_position_sheet1 = 6   # Target position in Sheet1
target_position_sheet2 = 6   # Target position in Sheet2

# Helper function to move a column within a worksheet
def move_column(sheet, column_letter, target_position):
    # Extract data from the column
    col_idx = column_index_from_string(column_letter)
    column_data = [cell[0] for cell in sheet.iter_rows(min_col=col_idx, max_col=col_idx, values_only=True)]

    # Delete the original column
    sheet.delete_cols(col_idx)

    # Insert the column at the target position
    sheet.insert_cols(target_position)
    for row_idx, value in enumerate(column_data, start=1):
        sheet.cell(row=row_idx, column=target_position, value=value)

# Move the specified columns in the respective sheets
move_column(sheet1, column_to_move_sheet1, target_position_sheet1)
move_column(sheet2, column_to_move_sheet2, target_position_sheet2)

# Save the workbook after modifying the sheets
new_wb.save('selected_columns_V1.xlsx')


# Step 2.2: create the docx file from Issue Description,Recommendation,Status Update
import pandas as pd
from docxtpl import DocxTemplate

# Load the template
template_path = 'template.docx'

# Define the path to the Excel file
excel_file_path = 'selected_columns_V1.xlsx'

# Read the 'High' and 'Medium_Low' worksheets from the Excel workbook
sheet1_df = pd.read_excel(excel_file_path, sheet_name='High')
sheet2_df = pd.read_excel(excel_file_path, sheet_name='Medium_Low')
sheet3_df = pd.read_excel(excel_file_path, sheet_name='Accepted')

# Define a function to create a docx file from a row
def create_docx_from_row(row, index, prefix):
    # Define the context for the current row
    context = {
        'Issue_Description': row['Issue Description'],
        'Recommendation': row['Recommendation'],
        'Status_Update': row['Status Update']
    }

    # Create a new DocxTemplate object for each iteration to avoid overwriting
    doc = DocxTemplate(template_path)

    # Render the template with the current context
    doc.render(context)

    # Save the document with a unique filename
    doc.save(f'{prefix}_generated_doc_{index + 1}.docx')

# Iterate over the rows in the 'High' DataFrame and create a docx file for each
for index, row in sheet1_df.iterrows():
    create_docx_from_row(row, index, 'High')

# Iterate over the rows in the 'Medium_Low' DataFrame and create a docx file for each
for index, row in sheet2_df.iterrows():
    create_docx_from_row(row, index, 'Medium_Low')

# Iterate over the rows in the 'Accepted' DataFrame and create a docx file for each
for index, row in sheet3_df.iterrows():
    create_docx_from_row(row, index, 'Accepted')

# Delete the three columns - 'Issue Description', 'Recommendation', 'Status Update'
# Clean column -  'Issue Description'

import pandas as pd

# Define the path to the Excel file
excel_file_path = 'selected_columns_V1.xlsx'
output_file_path = 'selected_columns_V1.xlsx'

#

# Read the three worksheets from the Excel workbook
sheet_high_df = pd.read_excel(excel_file_path, sheet_name='High')
sheet_medium_low_df = pd.read_excel(excel_file_path, sheet_name='Medium_Low')
sheet_accepted_df = pd.read_excel(excel_file_path, sheet_name='Accepted')

# Define the columns to delete
columns_to_delete = ['Issue Description', 'Recommendation', 'Status Update']

# Drop the specified columns from 'High','Medium_Low' and 'Accepted" worksheets
sheet_high_df = sheet_high_df.drop(columns=columns_to_delete)
sheet_medium_low_df = sheet_medium_low_df.drop(columns=columns_to_delete)
sheet_accepted_df = sheet_accepted_df.drop(columns=columns_to_delete)

# Create a new blank column 'Details in doc.' in both DataFrames
sheet_high_df['Details in doc'] = ''
sheet_medium_low_df['Details in doc'] = ''
sheet_accepted_df['Details in doc'] = ''

# Save the modified DataFrames back to an Excel file
with pd.ExcelWriter(output_file_path, engine='openpyxl') as writer:
    sheet_high_df.to_excel(writer, sheet_name='High', index=False)
    sheet_medium_low_df.to_excel(writer, sheet_name='Medium_Low', index=False)
    sheet_accepted_df.to_excel(writer, sheet_name='Accepted', index=False)


# Define the date format
import pandas as pd

# Define the path to the Excel file
excel_file_path = 'selected_columns_V1.xlsx'
output_file_path = 'selected_columns_V1.xlsx'

# Read the worksheets from the Excel workbook
sheet_high_df = pd.read_excel(excel_file_path, sheet_name='High')
sheet_medium_low_df = pd.read_excel(excel_file_path, sheet_name='Medium_Low')
sheet_accepted_df = pd.read_excel(excel_file_path, sheet_name='Accepted')

# Function to identify and format date columns
def format_date_columns(df):
    for column in df.columns:
        if pd.api.types.is_datetime64_any_dtype(df[column]):
            # If already datetime, format the dates
            df[column] = df[column].dt.strftime('%d-%b-%Y')
        elif pd.api.types.is_object_dtype(df[column]):
            try:
                # Attempt to convert to datetime and then format
                df[column] = pd.to_datetime(df[column], errors='raise')
                if pd.api.types.is_datetime64_any_dtype(df[column]):
                    df[column] = df[column].dt.strftime('%d-%b-%Y')
            except (ValueError, TypeError):
                # If conversion fails, skip formatting
                pass

# Apply the date formatting function to each DataFrame
format_date_columns(sheet_high_df)
format_date_columns(sheet_medium_low_df)
format_date_columns(sheet_accepted_df)

# Save the formatted DataFrames back to a new Excel file
with pd.ExcelWriter(output_file_path, engine='openpyxl') as writer:
    sheet_high_df.to_excel(writer, sheet_name='High', index=False)
    sheet_medium_low_df.to_excel(writer, sheet_name='Medium_Low', index=False)
    sheet_accepted_df.to_excel(writer, sheet_name='Accepted', index=False)


# format

from openpyxl import load_workbook
from openpyxl.styles import Border, Side
from openpyxl.styles import Alignment

# Define the border style
orange_border = Border(
    left=Side(style='thin', color='FFA500'),
    right=Side(style='thin', color='FFA500'),
    top=Side(style='thin', color='FFA500'),
    bottom=Side(style='thin', color='FFA500')
)

# Load the workbook to be modified
input_file_path = 'selected_columns_V1.xlsx'
output_file_path = 'selected_columns_V1.xlsx'
workbook = load_workbook(input_file_path)

# Define the border style (thin in this example) and color (orange)
orange_border = Border(
    left=Side(border_style="thin", color="FFA500"),
    right=Side(border_style="thin", color="FFA500"),
    top=Side(border_style="thin", color="FFA500"),
    bottom=Side(border_style="thin", color="FFA500")
)

# Apply the border style to all cells in all sheets
from openpyxl.styles import Alignment
for sheet in workbook.worksheets:
    for row in sheet.iter_rows():
        for cell in row:
            cell.border = orange_border

from openpyxl.styles import Font
# Define the font style (name, size, color, etc.)
font_style = Font(name="ING ME (body)", size=8, color="000000")  # Black color

# Apply the font style to all cells in the worksheet
for sheet in workbook.worksheets:
    for row in sheet.iter_rows():
        for cell in row:
            cell.font = font_style


# Adjust column widths:
# Iterate over each sheet in the workbook
for sheet in workbook.worksheets:
    # Adjust column widths and wrap text
    for column_cells in sheet.columns:
        max_length = 0
        column = column_cells[0].column_letter  # Get the column letter
        for cell in column_cells:
            # Wrap text in each cell
            cell.alignment = Alignment(vertical='top', wrap_text=True)

            # Determine the maximum length of the cell content in the column
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(cell.value)
            except:
                pass

        # Set the column width, adding a little extra space
    sheet.column_dimensions[column].width = max_length

# Save the workbook with the new styles
workbook.save(output_file_path)


# Step 3: ppt slides created by python-pptx and insert with docx files

# pip install pywin32

# pip install python-pptx
# Import necessary libraries
import openpyxl
import xlrd
from openpyxl import Workbook
from openpyxl.styles import NamedStyle
from datetime import datetime
import pandas as pd
from docxtpl import DocxTemplate
import win32com.client
import os
from pptx import Presentation
from pptx.util import Inches


# Step 3.1: create the ppt slides

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml

# Load the Excel file
file_path = 'selected_columns_V1.xlsx'
excel_data = pd.ExcelFile(file_path)

# Function to convert DataFrame to list of lists (including headers)
def convert_df_to_list(df):
    # Clean up text by removing line breaks, extra spaces, and "_x000D_" artifacts
    df = df.map(lambda x: x.replace('_x000D_', '').replace('\n', ' ').replace('\r', ' ').strip() if isinstance(x, str) else x)
    return [df.columns.tolist()] + df.fillna('').values.tolist()

# Load and convert each sheet's data, keeping the header row
sheets_data = {sheet: convert_df_to_list(excel_data.parse(sheet)) for sheet in excel_data.sheet_names}

# Function to estimate row height based on text content
def estimate_row_height(text, base_height=0.25, char_per_line=300):
    if isinstance(text, str):
        line_count = (len(text) // char_per_line) + 1
        return base_height * line_count
    return base_height

# Function to calculate the number of rows that can fit on a slide based on content
def calculate_rows_per_slide(data, row_height_func, max_height):
    rows = []
    current_height = 0
    for row in data[1:]:  # Skip header row
        row_height = max([row_height_func(cell) for cell in row])
        if current_height + row_height > max_height:
            break
        rows.append(row)
        current_height += row_height
    return len(rows)

# Function to split data into chunks based on the calculated rows per slide
def split_data_to_fit(data, rows_per_slide):
    return [data[:1] + data[i:i + rows_per_slide] for i in range(1, len(data), rows_per_slide)]

# Define the slide and table dimensions
slide_height = 8.0  # Total slide height in inches
title_height = 0.04  # Height reserved for the title in inches
margin_height = 4.0 # Additional margin height in inches
available_height = slide_height - title_height - margin_height

# Estimate rows per slide for each sheet
split_data = {}
for sheet, data in sheets_data.items():
    rows_per_slide = calculate_rows_per_slide(data, estimate_row_height, available_height)
    split_data[sheet] = split_data_to_fit(data, rows_per_slide)

def set_table_cell_border(cell, color):
    """Set the border of a table cell to the specified color."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for border in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        border_elm = parse_xml(
            f'<{border} w="12700" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill></{border}>')
        tcPr.append(border_elm)

def create_table(slide, title_text, data, column_widths):
    """Create a table on a slide with proper formatting."""
    # Set the slide title
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 98, 0)
    title.text_frame.paragraphs[0].font.size = Pt(25)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Set table position and size (adjusted to reduce space between title and table)
    left = Inches(0.3)
    top = Inches(0.5)
    width = Inches(10.0)
    height = prs.slide_height - top - Inches(0.5)  # Adjust height to maintain 0.5-inch distance from the bottom

    # Add the table
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.35) * rows).table

    # Set column widths
    for i, width in enumerate(column_widths):
        table.columns[i].width = width

    # Set header row
    for i, column_name in enumerate(data[0]):
        cell = table.cell(0, i)
        cell.text = column_name
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(8)
        p.font.name = 'ING Me (body)'
        p.font.color.rgb = RGBColor(0, 0, 0)  # Changed from orange to black
        p.alignment = PP_ALIGN.CENTER
        set_table_cell_border(cell, "FF6200")
        cell.fill.background()

    # Fill in the table rows
    for row_idx, row_data in enumerate(data[1:]):  # Start from the second row
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(7)
            p.font.name = 'ING Me (body)'
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT if col_idx > 0 else PP_ALIGN.CENTER
            set_table_cell_border(cell, "FF6200")
            cell.fill.background()

def add_slides_with_table(prs, title, data, column_widths):
    """Add slides with a table split over multiple slides if necessary."""
    for chunk in data:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        create_table(slide, title, chunk, column_widths)

def is_empty_slide(slide):
    """Check if a slide is empty (no significant shapes)."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            return False
        if shape.has_table:
            return False
    return True

def delete_empty_slides(prs):
    """Remove empty slides from the presentation."""
    slide_indexes_to_delete = []
    for i, slide in enumerate(prs.slides):
        if is_empty_slide(slide):
            slide_indexes_to_delete.append(i)

    for i in sorted(slide_indexes_to_delete, reverse=True):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

# Load the presentation
prs = Presentation()

# Define the column widths (adjust these based on the data)
high_columns_widths = [
    Inches(1.5), Inches(2.7), Inches(0.7), Inches(0.55), Inches(0.85), Inches(0.85), Inches(1.6), Inches(0.6)
]

medium_low_columns_widths = [
    Inches(1.5), Inches(2.7), Inches(0.7), Inches(0.55), Inches(0.85), Inches(0.85), Inches(1.7), Inches(0.6)
]

accepted_columns_widths = [
    Inches(1.5), Inches(2.5), Inches(0.7), Inches(0.85), Inches(0.85), Inches(0.85), Inches(1.4), Inches(0.6)
]

# Add slides with tables for each category
add_slides_with_table(prs, "High Findings Private Cloud (   in total)", split_data['High'], high_columns_widths)
add_slides_with_table(prs, "Medium Findings Private Cloud (   in total)", split_data['Medium_Low'], medium_low_columns_widths)
add_slides_with_table(prs, "Accepted risks Private Cloud (   in total)", split_data['Accepted'], accepted_columns_widths)

# Remove any empty slides
delete_empty_slides(prs)

# Save the presentation
prs.save('Cloud Compliance steerco.pptx')

# Step 3.2: insert the docs files

# Function to add a DOCX attachment in the "Details in doc" column
import win32com.client
import os

# Step 4: Insert DOCX files into the PowerPoint presentation
def insert_docx_as_icon_in_ppt(ppt_path, high_docx_files, medium_low_docx_files,accepted_docx_files, Details_in_doc_col_idx=7):
    App = win32com.client.Dispatch("PowerPoint.Application")
    App.Visible = True

    Presentation = App.Presentations.Open(os.path.abspath(ppt_path))

    docx_index_high = 0
    docx_index_medium_low = 0
    docx_index_accepted = 0

    for slide in Presentation.Slides:
        title_shape = slide.Shapes.Title
        if title_shape:
            title_text = title_shape.TextFrame.TextRange.Text
        else:
            title_text = ""

        for shape in slide.Shapes:
            if shape.HasTable:
                table = shape.Table

                for row_idx in range(2, table.Rows.Count + 1):
                    if "Medium Findings Private Cloud (   in total)" in title_text and docx_index_medium_low < len(medium_low_docx_files):
                        docx_file_path = medium_low_docx_files[docx_index_medium_low]
                        docx_index_medium_low += 1
                        cell = table.Cell(row_idx, Details_in_doc_col_idx)
                        left = cell.Shape.left + 10  # Adjusted position to move closer to the column
                        top = cell.Shape.Top + 10
                        width = 300  # Custom width for Medium  Findings
                        height = 45
                    elif "High Findings Private Cloud (   in total)" in title_text and docx_index_high < len(high_docx_files):
                        docx_file_path = high_docx_files[docx_index_high]
                        docx_index_high += 1
                        cell = table.Cell(row_idx, Details_in_doc_col_idx)
                        left = cell.Shape.left + 20  # Original position for High and Critical Findings
                        top = cell.Shape.Top + 10
                        width = 270  # Default width for High Findings
                        height = 45
                    elif "Accepted risks Private Cloud (   in total)" in title_text and docx_index_accepted < len(accepted_docx_files):
                        docx_file_path = accepted_docx_files[docx_index_accepted]
                        docx_index_accepted += 1
                        cell = table.Cell(row_idx, Details_in_doc_col_idx)
                        left = cell.Shape.left + 10  # Original position for Accepted
                        top = cell.Shape.Top + 10
                        width = 270  # Default width for Accepted Findings
                        height = 45
                    else:
                        continue

                    # Insert DOCX as an OLE object (icon)
                    ole_shape = slide.Shapes.AddOLEObject(
                        Left=left,
                        Top=top,
                        Width=width,
                        Height=height,
                        FileName=os.path.abspath(docx_file_path),
                        DisplayAsIcon=True,
                        IconFileName="C:\\Program Files\\Microsoft Office\\root\\Office16\\WINWORD.EXE",
                        IconLabel="Microsoft Word Document"
                    )

                    # Adjust the icon size using ScaleWidth and ScaleHeight instead of setting exact dimensions
                    ole_shape.LockAspectRatio = False
                    ole_shape.ScaleWidth(0.4, True)  # Scale the width by 0.4x
                    ole_shape.ScaleHeight(0.4, True)  # Scale the height by 0.4x

    Presentation.Save()
    # Presentation.Close()
    App.Quit()

# Define paths
ppt_path = 'Cloud Compliance steerco.pptx'
high_docx_files = [f"High_generated_doc_{i}.docx" for i in range(1, 25)]  # Adjust paths accordingly
medium_low_docx_files = [f"Medium_Low_generated_doc_{i}.docx" for i in range(1, 25)]  # Adjust paths accordingly
accepted_docx_files = [f"Accepted_generated_doc_{i}.docx" for i in range(1, 25)]  # Adjust paths accordingly

insert_docx_as_icon_in_ppt(ppt_path, high_docx_files, medium_low_docx_files, accepted_docx_files)

# Final Step: Adjust the pptx file to fit the screen

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

# Load the PowerPoint presentation
presentation = Presentation('Cloud Compliance steerco.pptx')


# Iterate through each slide
for slide in presentation.slides:
    # Find the title placeholder
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 0:  # Title placeholder
            # Set the title placeholder's position to the top of the slide
            shape.top = Inches(0.3)  # Adjust as needed to move the title to the top
            shape.left = Inches(0.3)  # Adjust as needed to center the title horizontally
            shape.width = presentation.slide_width - Inches(1)  # Adjust width to fit the slide

            # Set the title text alignment to center
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


# Save the modified presentation
presentation.save('Cloud Compliance steerco.pptx')



